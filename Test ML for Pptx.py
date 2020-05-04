"""
Created on Sun May  3 21:01:25 2020



@author: Marco Perdomo
"""
#!/usr/bin/env python
from pptx import Presentation
import glob
import pandas as pd
import PyPDF2 
import textract 

from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords


d = {'Title': [], 'Text': []}
for eachfile in glob.glob("*.pptx"):
    text_per_slide = []
    prs = Presentation(eachfile)
    d['Title'].append(eachfile)
    print("Extracting data from: %s" %eachfile)
    for slide in prs.slides:
        
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_per_slide.append(shape.text)
    d['Text'].append(str(text_per_slide))            

for eachfile in glob.glob("*.pdf"):
    d['Title'].append(eachfile)
    #Write a for-loop to open many files (leave a comment if you'd like to learn how).
    filename = '11 anatomy of wood(3).pdf'
    #open allows you to read the file.
    pdfFileObj = open(filename,'rb')
    #The pdfReader variable is a readable object that will be parsed.
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    #Discerning the number of pages will allow us to parse through all the pages.
    num_pages = pdfReader.numPages
    count = 0
    text = ""
    #The while loop will read each page.
    while count < num_pages:
        pageObj = pdfReader.getPage(count)
        count +=1
        text += pageObj.extractText()
    
    #This if statement exists to check if the above library returned words. It's done because PyPDF2 cannot read scanned files.
    if text != "":
       text = text
    #If the above returns as False, we run the OCR library textract to #convert scanned/image based PDF files into text.
    else:
       text = textract.process(fileurl, method='tesseract', language='eng')
       
    #Now we have a text variable that contains all the text derived from our PDF file. Type print(text) to see what it contains. It likely contains a lot of spaces, possibly junk such as '\n,' etc.
    #Now, we will clean our text variable and return it as a list of keywords.
    print(text)
    d['Text'].append(str(text))  


df = pd.DataFrame(d)
df.to_csv(r'SlidestoText.csv', index = False)


import keras
from keras.preprocessing.text import text_to_word_sequence
# define the document
test_text = 'Tu mama es mi perra'
# tokenize the document
result = text_to_word_sequence(test_text)
print(result)


